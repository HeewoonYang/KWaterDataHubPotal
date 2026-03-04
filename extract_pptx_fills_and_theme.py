#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Supplementary PPTX extractor: cell fills, table borders, and theme colors via XML.
"""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from lxml import etree
import zipfile

PPTX_PATH = r"D:\00_수공프로젝트\20260224_설계단계_산출물\작업중 산출물\(템플릿)KWDP-DH-DS-12-화면설계서-v0.x (원비율).pptx"

NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

def extract_cell_fill_from_xml(tc_elem):
    """Extract fill color from a table cell XML element."""
    tcPr = tc_elem.find('a:tcPr', NS)
    if tcPr is None:
        return "No tcPr"

    # Check solid fill
    solidFill = tcPr.find('a:solidFill', NS)
    if solidFill is not None:
        srgb = solidFill.find('a:srgbClr', NS)
        if srgb is not None:
            return f"SolidFill: #{srgb.get('val')}"
        schemeClr = solidFill.find('a:schemeClr', NS)
        if schemeClr is not None:
            return f"SolidFill: scheme={schemeClr.get('val')}"

    # Check no fill
    noFill = tcPr.find('a:noFill', NS)
    if noFill is not None:
        return "NoFill"

    return "Inherited (no fill specified)"

def extract_cell_borders_from_xml(tc_elem):
    """Extract border info from table cell XML."""
    tcPr = tc_elem.find('a:tcPr', NS)
    if tcPr is None:
        return {}

    borders = {}
    for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
        ln = tcPr.find(f'a:{border_name}', NS)
        if ln is not None:
            w = ln.get('w', 'default')
            solidFill = ln.find('a:solidFill', NS)
            color = "none"
            if solidFill is not None:
                srgb = solidFill.find('a:srgbClr', NS)
                if srgb is not None:
                    color = f"#{srgb.get('val')}"
                schemeClr = solidFill.find('a:schemeClr', NS)
                if schemeClr is not None:
                    color = f"scheme={schemeClr.get('val')}"
            noFill = ln.find('a:noFill', NS)
            if noFill is not None:
                color = "noFill"
            borders[border_name] = f"w={w} color={color}"
        else:
            borders[border_name] = "not specified"
    return borders


def main():
    prs = Presentation(PPTX_PATH)

    # ============================================================
    # EXTRACT THEME COLORS DIRECTLY FROM ZIP
    # ============================================================
    print("=" * 80)
    print("THEME EXTRACTION (from ZIP)")
    print("=" * 80)

    with zipfile.ZipFile(PPTX_PATH, 'r') as z:
        # List all files
        theme_files = [f for f in z.namelist() if 'theme' in f.lower()]
        print(f"  Theme files in archive: {theme_files}")

        for tf in theme_files:
            print(f"\n  --- {tf} ---")
            xml_content = z.read(tf)
            root = etree.fromstring(xml_content)

            # Color scheme
            for cs in root.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}clrScheme'):
                print(f"  Color Scheme: name='{cs.get('name')}'")
                for child in cs:
                    tag = child.tag.split('}')[-1]
                    srgb = child.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                    sysclr = child.find('{http://schemas.openxmlformats.org/drawingml/2006/main}sysClr')
                    if srgb is not None:
                        print(f"    {tag}: #{srgb.get('val')}")
                    elif sysclr is not None:
                        print(f"    {tag}: system={sysclr.get('val')} lastClr=#{sysclr.get('lastClr', 'N/A')}")

            # Font scheme
            for fs in root.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}fontScheme'):
                print(f"\n  Font Scheme: name='{fs.get('name')}'")
                majorFont = fs.find('{http://schemas.openxmlformats.org/drawingml/2006/main}majorFont')
                if majorFont is not None:
                    latin = majorFont.find('{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
                    ea = majorFont.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
                    print(f"    Major Font - Latin: {latin.get('typeface') if latin is not None else 'N/A'}")
                    print(f"    Major Font - EA: {ea.get('typeface') if ea is not None else 'N/A'}")
                    # Additional EA fonts
                    for font_elem in majorFont.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}font'):
                        print(f"    Major Font - {font_elem.get('script')}: {font_elem.get('typeface')}")

                minorFont = fs.find('{http://schemas.openxmlformats.org/drawingml/2006/main}minorFont')
                if minorFont is not None:
                    latin = minorFont.find('{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
                    ea = minorFont.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
                    print(f"    Minor Font - Latin: {latin.get('typeface') if latin is not None else 'N/A'}")
                    print(f"    Minor Font - EA: {ea.get('typeface') if ea is not None else 'N/A'}")
                    for font_elem in minorFont.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}font'):
                        print(f"    Minor Font - {font_elem.get('script')}: {font_elem.get('typeface')}")

    # ============================================================
    # EXTRACT TABLE CELL FILLS AND BORDERS FROM XML
    # ============================================================
    print("\n" + "=" * 80)
    print("TABLE CELL FILLS AND BORDERS (from XML)")
    print("=" * 80)

    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n  SLIDE {slide_idx + 1}:")

        for shape in slide.shapes:
            if shape.has_table:
                tbl_elem = shape._element.find('.//a:tbl', NS)
                if tbl_elem is None:
                    continue

                # Table grid / column widths from XML
                tblGrid = tbl_elem.find('a:tblGrid', NS)
                if tblGrid is not None:
                    gridCols = tblGrid.findall('a:gridCol', NS)
                    print(f"    Table Grid Columns (from XML):")
                    for ci, gc in enumerate(gridCols):
                        print(f"      Col[{ci}] w={gc.get('w')} EMU")

                rows = tbl_elem.findall('a:tr', NS)
                for row_idx, tr in enumerate(rows):
                    row_h = tr.get('h', 'N/A')
                    cells = tr.findall('a:tc', NS)
                    for col_idx, tc in enumerate(cells):
                        fill_info = extract_cell_fill_from_xml(tc)
                        borders = extract_cell_borders_from_xml(tc)

                        # Get merge info
                        gridSpan = tc.get('gridSpan', '1')
                        rowSpan = tc.get('rowSpan', '1')
                        hMerge = tc.get('hMerge', None)
                        vMerge = tc.get('vMerge', None)

                        # Only print cells that have content or special formatting
                        text_content = ''.join(t.text or '' for t in tc.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}t'))

                        merge_info = ""
                        if gridSpan != '1':
                            merge_info += f" gridSpan={gridSpan}"
                        if rowSpan != '1':
                            merge_info += f" rowSpan={rowSpan}"
                        if hMerge:
                            merge_info += f" hMerge={hMerge}"
                        if vMerge:
                            merge_info += f" vMerge={vMerge}"

                        # Print all cells for completeness
                        print(f"    Cell[{row_idx},{col_idx}]: fill={fill_info}{merge_info}")
                        if borders:
                            non_default = {k:v for k,v in borders.items() if v != "not specified"}
                            if non_default:
                                for bname, bval in non_default.items():
                                    print(f"      Border {bname}: {bval}")

    # ============================================================
    # EXTRACT LAYOUT LINE/SHAPE DETAILS FROM XML
    # ============================================================
    print("\n" + "=" * 80)
    print("LAYOUT SHAPE FILL COLORS (from XML)")
    print("=" * 80)

    for m_idx, master in enumerate(prs.slide_masters):
        for l_idx, layout in enumerate(master.slide_layouts):
            print(f"\n  Layout[{l_idx}]: '{layout.name}'")

            spTree = layout.element.find('.//p:cSld/p:spTree', NS)
            if spTree is None:
                continue

            for sp in spTree:
                tag = sp.tag.split('}')[-1]
                if tag not in ('sp', 'cxnSp', 'pic'):
                    continue

                # Get shape name
                nvSpPr = sp.find('.//p:nvSpPr/p:nvPr', NS)
                cNvPr = sp.find('.//p:nvSpPr/p:cNvPr', NS)
                if cNvPr is None:
                    cNvPr = sp.find('.//p:nvCxnSpPr/p:cNvPr', NS)

                name = cNvPr.get('name', 'N/A') if cNvPr is not None else 'N/A'

                # Get fill
                spPr = sp.find('.//p:spPr', NS)
                if spPr is None:
                    spPr = sp.find('p:spPr', NS)

                fill_str = "none"
                if spPr is not None:
                    solidFill = spPr.find('a:solidFill', NS)
                    if solidFill is not None:
                        srgb = solidFill.find('a:srgbClr', NS)
                        scheme = solidFill.find('a:schemeClr', NS)
                        if srgb is not None:
                            fill_str = f"#{srgb.get('val')}"
                        elif scheme is not None:
                            fill_str = f"scheme={scheme.get('val')}"

                    noFill = spPr.find('a:noFill', NS)
                    if noFill is not None:
                        fill_str = "noFill"

                    # Line
                    ln = spPr.find('a:ln', NS)
                    line_str = "none"
                    if ln is not None:
                        lw = ln.get('w', 'default')
                        lnSolidFill = ln.find('a:solidFill', NS)
                        lnColor = "none"
                        if lnSolidFill is not None:
                            lnSrgb = lnSolidFill.find('a:srgbClr', NS)
                            lnScheme = lnSolidFill.find('a:schemeClr', NS)
                            if lnSrgb is not None:
                                lnColor = f"#{lnSrgb.get('val')}"
                            elif lnScheme is not None:
                                lnColor = f"scheme={lnScheme.get('val')}"
                        line_str = f"w={lw} color={lnColor}"

                    print(f"    Shape '{name}': fill={fill_str} line={line_str}")


if __name__ == "__main__":
    main()
