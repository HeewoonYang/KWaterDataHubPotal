#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Comprehensive PPTX Design Specification Extractor
Extracts ALL design details from a PowerPoint template file.
"""

import sys
import os
import io

# Force UTF-8 output
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from lxml import etree

PPTX_PATH = r"D:\00_수공프로젝트\20260224_설계단계_산출물\작업중 산출물\(템플릿)KWDP-DH-DS-12-화면설계서-v0.x (원비율).pptx"

def emu_to_cm(emu):
    """Convert EMU to centimeters."""
    if emu is None:
        return None
    return round(emu / 914400 * 2.54, 4)

def emu_to_pt(emu):
    """Convert EMU to points."""
    if emu is None:
        return None
    return round(emu / 12700, 2)

def get_rgb_hex(color_obj):
    """Safely extract RGB hex from a color object."""
    try:
        if color_obj is None:
            return None
        if hasattr(color_obj, 'rgb') and color_obj.rgb is not None:
            return str(color_obj.rgb)
        if hasattr(color_obj, 'theme_color') and color_obj.theme_color is not None:
            tc = color_obj.theme_color
            brightness = getattr(color_obj, 'brightness', None)
            return f"Theme:{tc} (brightness={brightness})"
        return None
    except Exception as e:
        return f"Error:{e}"

def get_fill_info(fill):
    """Extract fill information."""
    try:
        if fill is None:
            return "None"
        fill_type = fill.type
        if fill_type is None:
            return "No fill / inherited"
        info = f"FillType={fill_type}"
        if hasattr(fill, 'fore_color') and fill.fore_color:
            rgb = get_rgb_hex(fill.fore_color)
            if rgb:
                info += f", ForeColor=#{rgb}"
        if hasattr(fill, 'back_color') and fill.back_color:
            rgb = get_rgb_hex(fill.back_color)
            if rgb:
                info += f", BackColor=#{rgb}"
        return info
    except Exception as e:
        return f"Fill error: {e}"

def get_alignment_str(alignment):
    """Convert alignment enum to string."""
    if alignment is None:
        return "None (inherited)"
    align_map = {
        PP_ALIGN.LEFT: "LEFT",
        PP_ALIGN.CENTER: "CENTER",
        PP_ALIGN.RIGHT: "RIGHT",
        PP_ALIGN.JUSTIFY: "JUSTIFY",
        PP_ALIGN.DISTRIBUTE: "DISTRIBUTE",
    }
    return align_map.get(alignment, str(alignment))

def print_paragraph_details(para, indent="        "):
    """Print detailed paragraph information."""
    print(f"{indent}Paragraph:")
    print(f"{indent}  Text: '{para.text}'")
    print(f"{indent}  Alignment: {get_alignment_str(para.alignment)}")

    # Paragraph-level spacing
    pf = para.paragraph_format
    if pf:
        print(f"{indent}  Space Before: {pf.space_before}")
        print(f"{indent}  Space After: {pf.space_after}")
        print(f"{indent}  Line Spacing: {pf.line_spacing}")
        print(f"{indent}  Level: {para.level}")

    for run_idx, run in enumerate(para.runs):
        print(f"{indent}  Run[{run_idx}]: '{run.text}'")
        font = run.font
        if font:
            print(f"{indent}    Font Name: {font.name}")
            print(f"{indent}    Font Size: {font.size} ({emu_to_pt(font.size) if font.size else 'None'} pt)")
            print(f"{indent}    Bold: {font.bold}")
            print(f"{indent}    Italic: {font.italic}")
            print(f"{indent}    Underline: {font.underline}")
            color_rgb = get_rgb_hex(font.color)
            print(f"{indent}    Color: {color_rgb}")
            if font.color and font.color.theme_color:
                print(f"{indent}    Theme Color: {font.color.theme_color}")
                if font.color.brightness is not None:
                    print(f"{indent}    Brightness: {font.color.brightness}")

def print_text_frame_details(tf, indent="      "):
    """Print text frame details."""
    print(f"{indent}Text Frame:")
    print(f"{indent}  Word Wrap: {tf.word_wrap}")
    print(f"{indent}  Auto Size: {tf.auto_size}")
    # Margins
    print(f"{indent}  Margin Left: {tf.margin_left} ({emu_to_cm(tf.margin_left)} cm)")
    print(f"{indent}  Margin Right: {tf.margin_right} ({emu_to_cm(tf.margin_right)} cm)")
    print(f"{indent}  Margin Top: {tf.margin_top} ({emu_to_cm(tf.margin_top)} cm)")
    print(f"{indent}  Margin Bottom: {tf.margin_bottom} ({emu_to_cm(tf.margin_bottom)} cm)")

    for p_idx, para in enumerate(tf.paragraphs):
        print(f"{indent}  Paragraph[{p_idx}]:")
        print_paragraph_details(para, indent + "    ")

def print_table_details(table, indent="      "):
    """Print table details."""
    print(f"{indent}Table: {table.rows.__len__()} rows x {len(table.columns)} cols")

    # Column widths
    for col_idx, col in enumerate(table.columns):
        print(f"{indent}  Column[{col_idx}] Width: {col.width} EMU ({emu_to_cm(col.width)} cm)")

    # Row heights
    for row_idx, row in enumerate(table.rows):
        print(f"{indent}  Row[{row_idx}] Height: {row.height} EMU ({emu_to_cm(row.height)} cm)")

    # Cell details
    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            print(f"{indent}  Cell[{row_idx},{col_idx}]:")
            print(f"{indent}    Text: '{cell.text}'")

            # Cell fill
            try:
                fill_info = get_fill_info(cell.fill)
                print(f"{indent}    Fill: {fill_info}")
            except:
                print(f"{indent}    Fill: (unable to read)")

            # Cell margins
            print(f"{indent}    Margin Left: {cell.margin_left}")
            print(f"{indent}    Margin Right: {cell.margin_right}")
            print(f"{indent}    Margin Top: {cell.margin_top}")
            print(f"{indent}    Margin Bottom: {cell.margin_bottom}")

            # Vertical anchor
            try:
                print(f"{indent}    Vertical Anchor: {cell.vertical_anchor}")
            except:
                pass

            # Text details in cell
            if cell.text_frame:
                for p_idx, para in enumerate(cell.text_frame.paragraphs):
                    print(f"{indent}    Para[{p_idx}]: '{para.text}' align={get_alignment_str(para.alignment)}")
                    for run_idx, run in enumerate(para.runs):
                        f = run.font
                        print(f"{indent}      Run[{run_idx}]: '{run.text}' font={f.name} size={f.size}({emu_to_pt(f.size) if f.size else None}pt) bold={f.bold} italic={f.italic} color={get_rgb_hex(f.color)}")

def get_border_info_from_xml(table_element):
    """Extract border info from table XML."""
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }
    try:
        tbl_pr = table_element.find('.//a:tblPr', nsmap)
        if tbl_pr is not None:
            print(f"      Table Properties XML: {etree.tostring(tbl_pr, pretty_print=True).decode()[:2000]}")
    except Exception as e:
        print(f"      Border XML extraction error: {e}")

def print_shape_details(shape, indent="    "):
    """Print comprehensive shape details."""
    print(f"{indent}Shape ID: {shape.shape_id}")
    print(f"{indent}Shape Name: '{shape.name}'")
    print(f"{indent}Shape Type: {shape.shape_type}")
    print(f"{indent}Position - Left: {shape.left} EMU ({emu_to_cm(shape.left)} cm)")
    print(f"{indent}Position - Top: {shape.top} EMU ({emu_to_cm(shape.top)} cm)")
    print(f"{indent}Size - Width: {shape.width} EMU ({emu_to_cm(shape.width)} cm)")
    print(f"{indent}Size - Height: {shape.height} EMU ({emu_to_cm(shape.height)} cm)")
    print(f"{indent}Rotation: {shape.rotation}")

    # Placeholder info
    if shape.is_placeholder:
        ph = shape.placeholder_format
        print(f"{indent}PLACEHOLDER:")
        print(f"{indent}  Placeholder idx: {ph.idx}")
        print(f"{indent}  Placeholder type: {ph.type}")

    # Fill info
    try:
        if hasattr(shape, 'fill'):
            fill_info = get_fill_info(shape.fill)
            print(f"{indent}Fill: {fill_info}")
    except Exception as e:
        print(f"{indent}Fill: Error - {e}")

    # Line/border info
    try:
        if hasattr(shape, 'line'):
            line = shape.line
            print(f"{indent}Line:")
            if line.fill and line.fill.type is not None:
                print(f"{indent}  Line Fill Type: {line.fill.type}")
                if hasattr(line, 'color') and line.color:
                    print(f"{indent}  Line Color: {get_rgb_hex(line.color)}")
            if line.width is not None:
                print(f"{indent}  Line Width: {line.width} EMU ({emu_to_pt(line.width)} pt)")
            if line.dash_style is not None:
                print(f"{indent}  Dash Style: {line.dash_style}")
    except Exception as e:
        print(f"{indent}Line: Error - {e}")

    # Text frame
    if shape.has_text_frame:
        print_text_frame_details(shape.text_frame, indent + "  ")

    # Table
    if shape.has_table:
        print_table_details(shape.table, indent + "  ")
        # Also extract border info from XML
        try:
            get_border_info_from_xml(shape._element)
        except Exception as e:
            print(f"{indent}  Border XML error: {e}")

    # Picture/Image
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = shape.image
            print(f"{indent}IMAGE:")
            print(f"{indent}  Content Type: {image.content_type}")
            print(f"{indent}  Size: {image.size}")
            print(f"{indent}  Filename: {image.filename if hasattr(image, 'filename') else 'N/A'}")
        except Exception as e:
            print(f"{indent}IMAGE Error: {e}")

    # Group shape
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        print(f"{indent}GROUP SHAPE - sub-shapes:")
        try:
            for sub_shape in shape.shapes:
                print(f"{indent}  --- Sub-Shape ---")
                print_shape_details(sub_shape, indent + "    ")
        except Exception as e:
            print(f"{indent}  Group iteration error: {e}")


def extract_theme_colors(prs):
    """Extract theme colors from the presentation."""
    print("\n" + "=" * 80)
    print("THEME COLORS")
    print("=" * 80)

    try:
        theme = prs.slide_masters[0].element
        nsmap = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        }

        # Try to find theme element
        theme_elements = theme.findall('.//a:theme', nsmap)
        clrScheme = theme.findall('.//a:clrScheme', nsmap)

        if clrScheme:
            for cs in clrScheme:
                print(f"  Color Scheme Name: {cs.get('name', 'N/A')}")
                for child in cs:
                    tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                    # Look for srgbClr or sysClr
                    srgb = child.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                    sysclr = child.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}sysClr')
                    if srgb is not None:
                        print(f"    {tag}: #{srgb.get('val', 'N/A')}")
                    elif sysclr is not None:
                        print(f"    {tag}: System={sysclr.get('val', 'N/A')} lastClr=#{sysclr.get('lastClr', 'N/A')}")
                    else:
                        print(f"    {tag}: {etree.tostring(child).decode()[:200]}")
        else:
            print("  No color scheme found in slide master element directly.")
            # Try via the theme part
            try:
                theme_part = prs.slide_masters[0].element
                # Print raw XML snippet for analysis
                xml_str = etree.tostring(theme_part, pretty_print=True).decode()
                # Find clrScheme in raw XML
                idx = xml_str.find('clrScheme')
                if idx > 0:
                    print(f"  Found clrScheme in XML at position {idx}")
                    snippet = xml_str[max(0, idx-50):idx+2000]
                    print(f"  Snippet:\n{snippet[:3000]}")
            except Exception as e2:
                print(f"  Theme extraction fallback error: {e2}")
    except Exception as e:
        print(f"  Theme color extraction error: {e}")

    # Also try getting theme from the presentation's theme part directly
    try:
        print("\n  --- Attempting direct theme part access ---")
        for rel in prs.slide_masters[0].part.rels.values():
            if 'theme' in rel.reltype:
                print(f"  Theme relationship found: {rel.reltype}")
                theme_part = rel.target_part
                theme_xml = etree.tostring(theme_part.element, pretty_print=True).decode()
                # Extract color scheme section
                start = theme_xml.find('<a:clrScheme')
                if start < 0:
                    start = theme_xml.find('clrScheme')
                if start >= 0:
                    end = theme_xml.find('</a:clrScheme>', start)
                    if end < 0:
                        end = start + 3000
                    print(f"  Theme Color Scheme XML:\n{theme_xml[start:end+20]}")

                # Extract font scheme
                fstart = theme_xml.find('<a:fontScheme')
                if fstart >= 0:
                    fend = theme_xml.find('</a:fontScheme>', fstart)
                    if fend < 0:
                        fend = fstart + 2000
                    print(f"\n  Theme Font Scheme XML:\n{theme_xml[fstart:fend+20]}")

                # Extract format scheme (fills, lines, effects)
                efstart = theme_xml.find('<a:fmtScheme')
                if efstart >= 0:
                    efend = theme_xml.find('</a:fmtScheme>', efstart)
                    if efend < 0:
                        efend = efstart + 3000
                    print(f"\n  Theme Format Scheme XML (first 3000 chars):\n{theme_xml[efstart:min(efend+20, efstart+3000)]}")
    except Exception as e:
        print(f"  Direct theme access error: {e}")


def extract_slide_background(slide, indent="  "):
    """Extract slide background information."""
    print(f"{indent}Background:")
    try:
        bg = slide.background
        if bg:
            fill = bg.fill
            fill_type = fill.type
            print(f"{indent}  Fill Type: {fill_type}")
            if fill_type is not None:
                try:
                    print(f"{indent}  Fore Color: {get_rgb_hex(fill.fore_color)}")
                except:
                    pass
                try:
                    print(f"{indent}  Back Color: {get_rgb_hex(fill.back_color)}")
                except:
                    pass
        else:
            print(f"{indent}  No background defined (inherited)")
    except Exception as e:
        print(f"{indent}  Background error: {e}")

    # Also check XML for background
    try:
        nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                 'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        bg_elem = slide.element.find('.//p:bg', nsmap)
        if bg_elem is not None:
            print(f"{indent}  Background XML: {etree.tostring(bg_elem, pretty_print=True).decode()[:1000]}")
        else:
            print(f"{indent}  No <p:bg> element found (background inherited from layout/master)")
    except Exception as e:
        print(f"{indent}  BG XML error: {e}")


def main():
    print(f"Opening: {PPTX_PATH}")
    print(f"File exists: {os.path.exists(PPTX_PATH)}")
    print(f"File size: {os.path.getsize(PPTX_PATH)} bytes")

    prs = Presentation(PPTX_PATH)

    # Slide dimensions
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    print(f"\n{'='*80}")
    print("PRESENTATION PROPERTIES")
    print(f"{'='*80}")
    print(f"Slide Width:  {slide_w} EMU = {emu_to_cm(slide_w)} cm = {round(slide_w/914400, 4)} inches")
    print(f"Slide Height: {slide_h} EMU = {emu_to_cm(slide_h)} cm = {round(slide_h/914400, 4)} inches")
    print(f"Aspect Ratio: {round(slide_w/slide_h, 4)}")

    # ==============================
    # SLIDE MASTERS
    # ==============================
    print(f"\n{'='*80}")
    print("SLIDE MASTERS")
    print(f"{'='*80}")
    for m_idx, master in enumerate(prs.slide_masters):
        print(f"\n  Master[{m_idx}]: name='{master.name if hasattr(master, 'name') else 'N/A'}'")
        # Master shapes
        for s_idx, shape in enumerate(master.shapes):
            print(f"    Master Shape[{s_idx}]: id={shape.shape_id} name='{shape.name}' type={shape.shape_type}")
            print(f"      Position: left={shape.left}({emu_to_cm(shape.left)}cm) top={shape.top}({emu_to_cm(shape.top)}cm)")
            print(f"      Size: w={shape.width}({emu_to_cm(shape.width)}cm) h={shape.height}({emu_to_cm(shape.height)}cm)")
            if shape.is_placeholder:
                ph = shape.placeholder_format
                print(f"      Placeholder: idx={ph.idx} type={ph.type}")
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        print(f"      Text: '{p.text}'")

        # Master background
        extract_slide_background(master, "    ")

    # ==============================
    # SLIDE LAYOUTS
    # ==============================
    print(f"\n{'='*80}")
    print("SLIDE LAYOUTS")
    print(f"{'='*80}")
    for m_idx, master in enumerate(prs.slide_masters):
        print(f"\n  From Master[{m_idx}]:")
        for l_idx, layout in enumerate(master.slide_layouts):
            print(f"\n    Layout[{l_idx}]: name='{layout.name}'")
            # Layout shapes
            for s_idx, shape in enumerate(layout.shapes):
                print(f"      Layout Shape[{s_idx}]: id={shape.shape_id} name='{shape.name}' type={shape.shape_type}")
                print(f"        Position: left={shape.left}({emu_to_cm(shape.left)}cm) top={shape.top}({emu_to_cm(shape.top)}cm)")
                print(f"        Size: w={shape.width}({emu_to_cm(shape.width)}cm) h={shape.height}({emu_to_cm(shape.height)}cm)")
                if shape.is_placeholder:
                    ph = shape.placeholder_format
                    print(f"        Placeholder: idx={ph.idx} type={ph.type}")
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text.strip():
                            print(f"        Text: '{p.text}'")
                            for run in p.runs:
                                f = run.font
                                print(f"          Run: '{run.text}' font={f.name} size={f.size}({emu_to_pt(f.size) if f.size else None}pt) bold={f.bold} color={get_rgb_hex(f.color)}")

    # ==============================
    # THEME COLORS
    # ==============================
    extract_theme_colors(prs)

    # ==============================
    # SLIDES
    # ==============================
    print(f"\n{'='*80}")
    print(f"SLIDES (Total: {len(prs.slides)})")
    print(f"{'='*80}")

    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n{'#'*80}")
        print(f"SLIDE {slide_idx + 1}")
        print(f"{'#'*80}")

        # Slide layout
        layout = slide.slide_layout
        print(f"  Layout Name: '{layout.name}'")

        # Slide dimensions (same for all slides)
        print(f"  Dimensions: {slide_w} x {slide_h} EMU = {emu_to_cm(slide_w)} x {emu_to_cm(slide_h)} cm")

        # Background
        extract_slide_background(slide, "  ")

        # Slide notes
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    print(f"  Notes: '{notes[:500]}'")
        except:
            pass

        # Shapes
        print(f"\n  Shapes (count: {len(slide.shapes)}):")
        for s_idx, shape in enumerate(slide.shapes):
            print(f"\n  --- Shape [{s_idx}] ---")
            print_shape_details(shape, "    ")

        print()  # Blank line between slides


if __name__ == "__main__":
    main()
