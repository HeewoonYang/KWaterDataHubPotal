# -*- coding: utf-8 -*-
"""
K-water DataHub Portal IA 재설계 PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== COLOR PALETTE =====
KWATER_BLUE = RGBColor(0x00, 0x5B, 0xAA)       # K-water 메인 블루
KWATER_DARK = RGBColor(0x1A, 0x27, 0x44)        # 다크 네이비
KWATER_LIGHT = RGBColor(0xE8, 0xF0, 0xFE)       # 라이트 블루
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
RED = RGBColor(0xF4, 0x43, 0x36)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
CYAN = RGBColor(0x00, 0x96, 0x88)

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

# ===== HELPER FUNCTIONS =====
def add_bg(slide, color=KWATER_DARK):
    """슬라이드 배경색 설정"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    """사각형 도형 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_rounded_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    """둥근 사각형 도형 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=12, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, default_size=11, default_color=BLACK, line_spacing=1.2):
    """여러 줄 텍스트 (리스트 형태) 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            text, size, bold, color = line_info, default_size, False, default_color
        else:
            text = line_info.get('text', '')
            size = line_info.get('size', default_size)
            bold = line_info.get('bold', False)
            color = line_info.get('color', default_color)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = '맑은 고딕'
        p.space_after = Pt(2)
    return txBox

def add_card(slide, left, top, width, height, title, items, title_color=KWATER_BLUE, bg_color=WHITE):
    """카드 UI 컴포넌트"""
    # 카드 배경
    card = add_rounded_shape(slide, left, top, width, height, fill_color=bg_color, border_color=RGBColor(0xE0, 0xE0, 0xE0))
    # 타이틀 바
    add_shape(slide, left + Inches(0.02), top + Inches(0.02), width - Inches(0.04), Inches(0.4), fill_color=title_color)
    add_text(slide, left + Inches(0.15), top + Inches(0.06), width - Inches(0.3), Inches(0.35), title, font_size=11, bold=True, color=WHITE)
    # 아이템
    y_offset = top + Inches(0.5)
    for item in items:
        add_text(slide, left + Inches(0.15), y_offset, width - Inches(0.3), Inches(0.22), item, font_size=9, color=GRAY)
        y_offset += Inches(0.22)

def add_slide_number(slide, num, total):
    """슬라이드 번호"""
    add_text(slide, Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3), f'{num}/{total}', font_size=8, color=RGBColor(0x99,0x99,0x99), alignment=PP_ALIGN.RIGHT)

def add_header_bar(slide):
    """공통 상단 헤더"""
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), fill_color=KWATER_BLUE)

TOTAL_SLIDES = 14

# ======================================================================
# SLIDE 1: 표지
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, KWATER_DARK)

# 로고 영역
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), fill_color=KWATER_BLUE)

# 메인 타이틀
add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.7),
         'K-water DataHub Portal', font_size=42, bold=True, color=WHITE)
add_text(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.6),
         'IA (Information Architecture) 재설계', font_size=28, bold=False, color=RGBColor(0x4F, 0xC3, 0xF7))
add_text(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.4),
         '한국수자원공사 디지털플랫폼 구축사업 — 데이터허브 포털', font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA))

# 구분선
add_shape(slide, Inches(1), Inches(4.0), Inches(4), Inches(0.03), fill_color=KWATER_BLUE)

# 버전/날짜
info_lines = [
    {'text': 'v3.0  |  설계 단계  |  UI/UX 프로토타입', 'size': 13, 'color': RGBColor(0xBB,0xBB,0xBB)},
    {'text': '', 'size': 8},
    {'text': '2026.02', 'size': 13, 'color': RGBColor(0xBB,0xBB,0xBB)},
]
add_multiline_text(slide, Inches(1), Inches(4.3), Inches(6), Inches(1.5), info_lines)

# 우측 키워드 박스
keywords = [
    '계측DB 4개  (RWIS · HDAPS · GIOS · 스마트미터링)',
    '자산DB 6개  (자산 · 시설 · 수질 · 운영 · 관망통합 · 기상청)',
    'Kafka 실시간 스트리밍  |  CDC 변경데이터 수집',
    'LLM/AI 질의응답  |  그래프DB 온톨로지',
    'DT 5종 API 데이터 연계 제공',
]
y = Inches(4.3)
for kw in keywords:
    s = add_rounded_shape(slide, Inches(7.5), y, Inches(5.3), Inches(0.35), fill_color=RGBColor(0x25, 0x35, 0x55))
    s.text_frame.paragraphs[0].text = kw
    s.text_frame.paragraphs[0].font.size = Pt(10)
    s.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
    s.text_frame.paragraphs[0].font.name = '맑은 고딕'
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    y += Inches(0.42)

add_slide_number(slide, 1, TOTAL_SLIDES)


# ======================================================================
# SLIDE 2: 목차
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.5),
         'CONTENTS', font_size=28, bold=True, color=KWATER_DARK)
add_shape(slide, Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04), fill_color=KWATER_BLUE)

toc_items = [
    ('01', '프로젝트 개요 및 아키텍처'),
    ('02', '원천시스템 현황 (계측DB 4 + 자산DB 6)'),
    ('03', 'IA 전체 구조 (Top Nav + Sidebar)'),
    ('04', '① 통합대시보드 상세'),
    ('05', '② 카탈로그·탐색 상세'),
    ('06', '③ 수집·정제 상세'),
    ('07', '④ 유통·활용 상세'),
    ('08', '⑤ 메타데이터 관리 상세'),
    ('09', '⑥ 시스템관리 상세'),
    ('10', '역할별 접근권한 매트릭스 (RBAC)'),
    ('11', '화면 플로우 / 네비게이션 맵'),
    ('12', '현재 대비 변경점 요약'),
    ('13', '향후 일정 및 Next Steps'),
]

y = Inches(1.5)
for num, title in toc_items:
    # 번호
    s = add_rounded_shape(slide, Inches(1.2), y, Inches(0.55), Inches(0.38), fill_color=KWATER_BLUE)
    s.text_frame.paragraphs[0].text = num
    s.text_frame.paragraphs[0].font.size = Pt(12)
    s.text_frame.paragraphs[0].font.bold = True
    s.text_frame.paragraphs[0].font.color.rgb = WHITE
    s.text_frame.paragraphs[0].font.name = '맑은 고딕'
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # 타이틀
    add_text(slide, Inches(1.95), y + Inches(0.04), Inches(8), Inches(0.35), title, font_size=13, color=KWATER_DARK)
    y += Inches(0.43)

add_slide_number(slide, 2, TOTAL_SLIDES)


# ======================================================================
# SLIDE 3: 프로젝트 개요 및 아키텍처
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.4),
         '01  프로젝트 개요 및 시스템 아키텍처', font_size=22, bold=True, color=KWATER_DARK)
add_shape(slide, Inches(0.8), Inches(0.75), Inches(2), Inches(0.04), fill_color=KWATER_BLUE)

# 아키텍처 다이어그램 (텍스트 기반)
layers = [
    {'label': '사용자', 'items': '모바일 사용자  |  웹 사용자  |  웹 관리자', 'color': RGBColor(0x26, 0x32, 0x38), 'y': 1.2},
    {'label': '서비스계', 'items': 'WEB서버 VM(x4)  |  WAS서버 VM(x4)  |  모니터링서버 VM(x4)', 'color': KWATER_BLUE, 'y': 2.0},
    {'label': '분석계', 'items': 'ODSDB PM(x3)  |  분석DB PM(x4)', 'color': RGBColor(0x00, 0x96, 0x88), 'y': 2.8},
    {'label': '온톨로지/AI', 'items': 'LLM GM(x1) GPU:B300  |  그래프DB PM(x1)  |  마인즈DB VM(x4)', 'color': PURPLE, 'y': 3.6},
    {'label': '연계', 'items': 'MCP서버 VM(x4)  |  API서버 VM(x4)  →  DT 5종 (정수장/관로/댐/발전/스마트)', 'color': ORANGE, 'y': 4.4},
    {'label': '수집계', 'items': '수집DB PM(x3)  |  카프카 VM(x6)  |  CDC서버 VM(x4)  |  수집Agent VM(x5)  |  프로듀서 VM(x4)', 'color': RGBColor(0xC6, 0x28, 0x28), 'y': 5.2},
]

for layer in layers:
    ly = Inches(layer['y'])
    # 라벨
    s = add_rounded_shape(slide, Inches(0.8), ly, Inches(1.5), Inches(0.55), fill_color=layer['color'])
    s.text_frame.paragraphs[0].text = layer['label']
    s.text_frame.paragraphs[0].font.size = Pt(10)
    s.text_frame.paragraphs[0].font.bold = True
    s.text_frame.paragraphs[0].font.color.rgb = WHITE
    s.text_frame.paragraphs[0].font.name = '맑은 고딕'
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # 내용
    s2 = add_rounded_shape(slide, Inches(2.5), ly, Inches(7.5), Inches(0.55), fill_color=LIGHT_GRAY, border_color=RGBColor(0xDD,0xDD,0xDD))
    s2.text_frame.paragraphs[0].text = layer['items']
    s2.text_frame.paragraphs[0].font.size = Pt(9)
    s2.text_frame.paragraphs[0].font.color.rgb = KWATER_DARK
    s2.text_frame.paragraphs[0].font.name = '맑은 고딕'
    s2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 우측: 원천시스템 박스
add_text(slide, Inches(10.5), Inches(1.0), Inches(2.5), Inches(0.35),
         '원천시스템 (10개 DB)', font_size=12, bold=True, color=KWATER_DARK)
add_shape(slide, Inches(10.5), Inches(1.35), Inches(1.5), Inches(0.03), fill_color=KWATER_BLUE)

# 계측DB
add_text(slide, Inches(10.5), Inches(1.55), Inches(2.5), Inches(0.3),
         '계측DB (4개) — 실시간', font_size=10, bold=True, color=RGBColor(0xC6,0x28,0x28))
meas_dbs = ['RWIS (원격수위정보)', 'HDAPS (수력발전자동화)', 'GIOS (지하수관측)', '스마트미터링']
y = Inches(1.9)
for db in meas_dbs:
    s = add_rounded_shape(slide, Inches(10.5), y, Inches(2.5), Inches(0.3), fill_color=RGBColor(0xFF, 0xEB, 0xEE))
    s.text_frame.paragraphs[0].text = db
    s.text_frame.paragraphs[0].font.size = Pt(8)
    s.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xC6,0x28,0x28)
    s.text_frame.paragraphs[0].font.name = '맑은 고딕'
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    y += Inches(0.35)

# 자산DB
add_text(slide, Inches(10.5), Inches(3.4), Inches(2.5), Inches(0.3),
         '자산DB (6개) — CDC/배치', font_size=10, bold=True, color=KWATER_BLUE)
asset_dbs = ['자산DB (SAP)', '시설DB (SAP)', '수질DB (ORACLE)', '운영DB (SAP)', '관망통합DB (ORACLE)', '기상청DB (외부 API)']
y = Inches(3.75)
for db in asset_dbs:
    s = add_rounded_shape(slide, Inches(10.5), y, Inches(2.5), Inches(0.3), fill_color=KWATER_LIGHT)
    s.text_frame.paragraphs[0].text = db
    s.text_frame.paragraphs[0].font.size = Pt(8)
    s.text_frame.paragraphs[0].font.color.rgb = KWATER_BLUE
    s.text_frame.paragraphs[0].font.name = '맑은 고딕'
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    y += Inches(0.35)

add_slide_number(slide, 3, TOTAL_SLIDES)


# ======================================================================
# SLIDE 4: 원천시스템 현황
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.4),
         '02  원천시스템 현황 — 계측DB 4개 + 자산DB 6개', font_size=22, bold=True, color=KWATER_DARK)
add_shape(slide, Inches(0.8), Inches(0.75), Inches(2), Inches(0.04), fill_color=KWATER_BLUE)

# 계측DB 섹션
add_text(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.35),
         '실시간 계측DB (4개)', font_size=16, bold=True, color=RGBColor(0xC6,0x28,0x28))
add_text(slide, Inches(0.8), Inches(1.45), Inches(5), Inches(0.25),
         '단방향전송장치 / 프로듀서 VM(x4) → 카프카 VM(x6) → 수집DB PM(x3)', font_size=9, color=GRAY)

meas_cards = [
    ('RWIS', '원격수위정보시스템', ['전국 관측소 수위 데이터', '실시간 CDC 수집', '데이터량: 1,234,567건/일', '수집주기: 상시 (1분 간격)']),
    ('HDAPS', '수력발전자동화시스템', ['수력발전소 운영 데이터', '발전량/수위/유량 통합', '데이터량: 856,000건/일', '수집주기: 상시']),
    ('GIOS', '지하수관측시스템', ['전국 지하수 관정 데이터', '지하수위/수온/수질', '데이터량: 432,000건/일', '수집주기: 10분 간격']),
    ('스마트미터링', '원격검침시스템', ['스마트 수도미터 검침', '가정/산업용 사용량', '데이터량: 2,800,000건/일', '수집주기: 상시 (실시간)']),
]

x = Inches(0.8)
for name, subtitle, items in meas_cards:
    add_card(slide, x, Inches(1.8), Inches(2.9), Inches(2.0), f'{name} — {subtitle}', items, title_color=RGBColor(0xC6,0x28,0x28))
    x += Inches(3.05)

# 자산DB 섹션
add_text(slide, Inches(0.8), Inches(4.1), Inches(5), Inches(0.35),
         '운영 자산DB (6개)', font_size=16, bold=True, color=KWATER_BLUE)
add_text(slide, Inches(0.8), Inches(4.45), Inches(6), Inches(0.25),
         'CDC서버 VM(x4) / 수집Agent VM(x5) → 수집DB PM(x3)', font_size=9, color=GRAY)

asset_cards = [
    ('자산DB', 'SAP', ['시설물 자산 마스터', '취득/처분/감가상각', 'CDC 배치수집']),
    ('시설DB', 'SAP', ['정수장/관로/펌프장', '시설물 제원/점검이력', 'CDC 배치수집']),
    ('수질DB', 'ORACLE', ['수질검사 결과', 'pH/탁도/잔류염소', 'CDC 배치수집']),
    ('운영DB', 'SAP', ['일일 운영실적', '생산량/송수량/배수량', 'CDC 배치수집']),
    ('관망통합DB', 'ORACLE', ['관로 네트워크', '배관/밸브/소화전', 'CDC 배치수집']),
    ('기상청DB', '외부 API', ['기상 관측 데이터', '온도/강수/풍속', 'REST API 연계']),
]

x = Inches(0.8)
for name, src, items in asset_cards:
    add_card(slide, x, Inches(4.8), Inches(1.95), Inches(1.55), f'{name} ({src})', items, title_color=KWATER_BLUE)
    x += Inches(2.05)

add_slide_number(slide, 4, TOTAL_SLIDES)


# ======================================================================
# SLIDE 5: IA 전체 구조
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.4),
         '03  IA 전체 구조 — Top Navigation + Sidebar 메뉴', font_size=22, bold=True, color=KWATER_DARK)
add_shape(slide, Inches(0.8), Inches(0.75), Inches(2), Inches(0.04), fill_color=KWATER_BLUE)

# Top Nav 시뮬레이션
add_shape(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5), fill_color=KWATER_DARK)
nav_items = ['DataHub Portal', '①통합대시보드', '②카탈로그·탐색', '③수집·정제', '④유통·활용', '⑤메타데이터', '⑥시스템관리']
nav_x = Inches(0.7)
for i, item in enumerate(nav_items):
    w = Inches(1.5) if i == 0 else Inches(1.45)
    add_text(slide, nav_x, Inches(1.15), w, Inches(0.4), item,
             font_size=10 if i > 0 else 11, bold=(i==0), color=WHITE if i != 1 else RGBColor(0x4F,0xC3,0xF7))
    nav_x += w

# 6개 대메뉴 카드
menus = [
    ('① 통합대시보드', KWATER_BLUE, [
        '1.1 업무중심홈 (개인화)',
        '1.2 데이터 워크플로우',
        '1.3 데이터 품질현황 종합',
        '1.4 원천시스템 모니터링',
        '    ├ 계측DB 대시보드 (4개)',
        '    └ 자산DB 대시보드 (6개)',
        '1.5 AI 검색 및 질의응답',
        '1.6 시각화·차트 갤러리',
        '1.7 메인화면(위젯) 설정',
    ]),
    ('② 카탈로그·탐색', RGBColor(0x00,0x96,0x88), [
        '2.1 통합데이터 검색',
        '2.2 데이터셋 상세정보',
        '2.3 지식그래프·온톨로지',
        '2.4 데이터 리니지(계보)',
        '2.5 내 보관함',
        '2.6 데이터 요청·피드백',
    ]),
    ('③ 수집·정제', RGBColor(0xC6,0x28,0x28), [
        '3.1 파이프라인 관리',
        '3.2 파이프라인 등록',
        '3.3 원천시스템 연계관리',
        '    ├ CDC 연계현황',
        '    ├ Kafka 스트리밍',
        '    ├ 외부연계 (기상청/제타7)',
        '    └ 단방향전송 모니터링',
        '3.4 연계아키텍처 구조도',
        '3.5 수집현황·모니터링',
        '3.6 실시간 로그·알림',
        '3.7 정제·변환 관리 (dbt)',
    ]),
    ('④ 유통·활용', ORANGE, [
        '4.1 Data Product 구성·관리',
        '4.2 비식별화·가명처리',
        '4.3 데이터 신청·승인',
        '4.4 유통 API 관리',
        '4.5 외부시스템 데이터 제공',
        '4.6 데이터 유통통계',
    ]),
    ('⑤ 메타데이터 관리', PURPLE, [
        '5.1 표준용어·코드관리',
        '5.2 분류체계·태그관리',
        '5.3 데이터모델관리',
        '5.4 DQ 검증·프로파일링',
        '5.5 온톨로지 관리',
    ]),
    ('⑥ 시스템관리', RGBColor(0x45,0x5A,0x64), [
        '6.1 조직 및 사용자관리',
        '6.2 권한 및 역할관리 (RBAC)',
        '6.3 데이터등급·보안정책',
        '6.4 인프라 현황',
        '    ├ 수집계 / 분석계 / 서비스계',
        '6.5 LLM·AI 엔진관리',
        '6.6 연계인터페이스 모니터링',
        '6.7 접속통계·감사로그',
    ]),
]

x = Inches(0.5)
card_w = Inches(2.0)
for title, color, items in menus:
    # 타이틀
    s = add_rounded_shape(slide, x, Inches(1.85), card_w, Inches(0.38), fill_color=color)
    s.text_frame.paragraphs[0].text = title
    s.text_frame.paragraphs[0].font.size = Pt(10)
    s.text_frame.paragraphs[0].font.bold = True
    s.text_frame.paragraphs[0].font.color.rgb = WHITE
    s.text_frame.paragraphs[0].font.name = '맑은 고딕'
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # 아이템 목록
    y = Inches(2.35)
    for item in items:
        is_sub = item.startswith('    ')
        add_text(slide, x + Inches(0.05), y, card_w - Inches(0.1), Inches(0.2),
                 item, font_size=7.5 if is_sub else 8.5, color=GRAY if is_sub else KWATER_DARK)
        y += Inches(0.22)
    x += Inches(2.08)

# 하단 요약
add_shape(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8), fill_color=LIGHT_GRAY)
summary_text = '총 38개 화면 + 14개 서브뷰  |  Top Nav 6개 대메뉴  |  Sidebar 2Depth 아코디언  |  역할별 RBAC 메뉴 필터링'
add_text(slide, Inches(0.8), Inches(6.45), Inches(11.8), Inches(0.5), summary_text, font_size=12, bold=True, color=KWATER_DARK, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 5, TOTAL_SLIDES)


# ======================================================================
# SLIDE 6: ① 통합대시보드 상세
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.4),
         '04  ① 통합대시보드 상세', font_size=22, bold=True, color=KWATER_DARK)
add_shape(slide, Inches(0.8), Inches(0.75), Inches(2), Inches(0.04), fill_color=KWATER_BLUE)

dash_screens = [
    ('1.1 업무중심홈', KWATER_BLUE, [
        'AI 퀵서치 바',
        'KPI 4종 카드',
        '나의 업무현황',
        '품질 요약 (4대 지표)',
        '최근활동 / 공지',
        '★ 원천시스템 10개 DB 미니 상태 [신규]',
    ]),
    ('1.2 업무프로세스', RGBColor(0x00,0x96,0x88), [
        '5단계 프로세스 플로우',
        '(탐색→신청→승인→활용→환류)',
        '단계별 업무 목록 테이블',
        '워크플로우 상세 드로어',
    ]),
    ('1.3 품질현황 종합', GREEN, [
        '전사 품질 점수',
        '(완전성/정합성/유효성/적시성)',
        '★ 원천시스템별 품질 점수 [신규]',
        '  → 계측DB 4 + 자산DB 6',
        '품질 추이 차트',
        '품질 이슈 목록 / 조치현황',
    ]),
    ('1.5 AI 검색·질의응답', PURPLE, [
        '대화형 채팅 UI',
        '사고과정(Thinking) 시각화',
        '인라인 차트/테이블 응답',
        '역할별 추천 질문',
        '★ LLM 엔진상태 표시 [신규]',
    ]),
]

x = Inches(0.8)
for title, color, items in dash_screens:
    add_card(slide, x, Inches(1.2), Inches(2.9), Inches(2.5), title, items, title_color=color)
    x += Inches(3.05)

# 1.4 원천시스템 모니터링 (핵심 - 크게)
add_text(slide, Inches(0.8), Inches(3.95), Inches(10), Inches(0.35),
         '1.4 원천시스템 모니터링 (핵심 신규 화면)', font_size=14, bold=True, color=RGBColor(0xC6,0x28,0x28))

# 계측DB
add_shape(slide, Inches(0.8), Inches(4.4), Inches(5.8), Inches(2.5), fill_color=RGBColor(0xFF,0xF8,0xF0), border_color=RGBColor(0xC6,0x28,0x28))
add_text(slide, Inches(1.0), Inches(4.5), Inches(5), Inches(0.3),
         '1.4.1 계측DB 대시보드 (4개 DB)', font_size=12, bold=True, color=RGBColor(0xC6,0x28,0x28))

meas_items = [
    ('RWIS', '관측소 맵 / 시계열차트 / 임계치 현황'),
    ('HDAPS', '발전소별 현황 / 발전량 추이'),
    ('GIOS [신규]', '관정별 지하수위 / 수질 현황'),
    ('스마트미터링', '검침현황 / 사용량분석 / 이상탐지'),
]
y = Inches(4.85)
for name, desc in meas_items:
    s = add_rounded_shape(slide, Inches(1.0), y, Inches(1.5), Inches(0.35), fill_color=RGBColor(0xC6,0x28,0x28))
    s.text_frame.paragraphs[0].text = name
    s.text_frame.paragraphs[0].font.size = Pt(9)
    s.text_frame.paragraphs[0].font.bold = True
    s.text_frame.paragraphs[0].font.color.rgb = WHITE
    s.text_frame.paragraphs[0].font.name = '맑은 고딕'
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(slide, Inches(2.65), y + Inches(0.05), Inches(3.7), Inches(0.3), desc, font_size=9, color=GRAY)
    y += Inches(0.45)

# 자산DB
add_shape(slide, Inches(6.8), Inches(4.4), Inches(5.8), Inches(2.5), fill_color=RGBColor(0xF0,0xF4,0xFF), border_color=KWATER_BLUE)
add_text(slide, Inches(7.0), Inches(4.5), Inches(5), Inches(0.3),
         '1.4.2 자산DB 대시보드 (6개 DB) [완전 신규]', font_size=12, bold=True, color=KWATER_BLUE)

asset_items = [
    ('자산DB', '자산 등록/변경 현황'),
    ('시설DB', '시설물 현황/점검이력'),
    ('수질DB', '수질검사 결과 현황'),
    ('운영DB', '운영 실적 현황'),
    ('관망통합DB', '관로/배관 현황'),
    ('기상청DB [신규]', '기상 관측 데이터'),
]
y = Inches(4.85)
for name, desc in asset_items:
    s = add_rounded_shape(slide, Inches(7.0), y, Inches(1.5), Inches(0.28), fill_color=KWATER_BLUE)
    s.text_frame.paragraphs[0].text = name
    s.text_frame.paragraphs[0].font.size = Pt(8)
    s.text_frame.paragraphs[0].font.bold = True
    s.text_frame.paragraphs[0].font.color.rgb = WHITE
    s.text_frame.paragraphs[0].font.name = '맑은 고딕'
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(slide, Inches(8.65), y + Inches(0.02), Inches(3.7), Inches(0.25), desc, font_size=9, color=GRAY)
    y += Inches(0.35)

add_slide_number(slide, 6, TOTAL_SLIDES)


# ======================================================================
# SLIDE 7: ② 카탈로그·탐색 상세
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.4),
         '05  ② 카탈로그·탐색 상세', font_size=22, bold=True, color=KWATER_DARK)
add_shape(slide, Inches(0.8), Inches(0.75), Inches(2), Inches(0.04), fill_color=RGBColor(0x00,0x96,0x88))

cat_screens = [
    ('2.1 통합데이터 검색', [
        '통합 검색바 (도메인/태그/표준용어)',
        '★ 원천시스템 필터 10개 DB 명시 [개선]',
        '  계측: RWIS | HDAPS | GIOS | 스마트미터링',
        '  자산: 자산 | 시설 | 수질 | 운영 | 관망통합 | 기상청',
        '필터 (보안등급/수집방식/데이터유형)',
        '태그 필터 (도메인별 태그 칩)',
        '검색결과 카드 목록 + 정렬',
    ]),
    ('2.2 데이터셋 상세정보', [
        '기본정보 (명칭/설명/원천/수집방식/등급)',
        '스키마·컬럼 정보 탭',
        '샘플 데이터 미리보기 탭',
        '품질 정보 탭',
        '변경이력 탭',
        '★ 원천시스템 정보 표시 [신규]',
        '액션: 보관함/리니지/신청/API',
    ]),
    ('2.3 지식그래프·온톨로지', [
        '전사 데이터 관계 그래프',
        '노드 검색/필터',
        '노드 클릭 상세 패널',
        '★ 그래프DB(Neo4j) 기반 시각화',
    ]),
    ('2.4 데이터 리니지(계보)', [
        '리니지 대상 선택',
        '계층별 뷰',
        '  Source→Staging→Mart→Serving',
        '★ 전사 데이터 흐름도 [신규]',
        '  계측DB(4) → 카프카/CDC',
        '  → 수집DB → ODS → 분석DB → API/DT',
        '노드 상세 사이드 패널',
        '영향분석 모드',
    ]),
    ('2.5 내 보관함', [
        '즐겨찾기 데이터셋 목록',
        '최근 조회 이력',
        '구독 알림 설정',
    ]),
    ('2.6 데이터 요청·피드백', [
        '신규 데이터 요청 폼',
        '요청 현황 / 처리 이력',
        '피드백 (평점/의견)',
    ]),
]

x = Inches(0.5)
for title, items in cat_screens:
    h = max(Inches(2.8), Inches(0.5 + len(items) * 0.22))
    add_card(slide, x, Inches(1.2), Inches(2.0), h, title, items, title_color=RGBColor(0x00,0x96,0x88))
    x += Inches(2.08)

add_slide_number(slide, 7, TOTAL_SLIDES)


# ======================================================================
# SLIDE 8: ③ 수집·정제 상세
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.4),
         '06  ③ 수집·정제 상세', font_size=22, bold=True, color=KWATER_DARK)
add_shape(slide, Inches(0.8), Inches(0.75), Inches(2), Inches(0.04), fill_color=RGBColor(0xC6,0x28,0x28))

# 상단: 기존 메뉴
col_top = [
    ('3.1 파이프라인 관리', ['전체 파이프라인 목록', '필터 (수집방식/상태/원천)', '상태별 카운트', '상세 드릴다운']),
    ('3.2 파이프라인 등록', ['Step1: 원천시스템 선택 (10개 DB)', 'Step2: 수집방식', 'Step3: 스케줄/트리거', 'Step4: 타겟 (수집→ODS→분석) [신규]', 'Step5: 검증 및 등록']),
    ('3.4 아키텍처 구조도', ['전체 아키텍처 인터랙티브 [개선]', '실시간 상태 오버레이', '계층별 뷰']),
    ('3.5 수집현황 모니터링', ['전체 수집 통계', '★ 원천시스템별 수집량 (10개)', '★ 수집DB PM(x3) 용량/성능', '오류/지연 알림 목록']),
]

x = Inches(0.5)
for title, items in col_top:
    add_card(slide, x, Inches(1.1), Inches(3.0), Inches(2.0), title, items, title_color=RGBColor(0xC6,0x28,0x28))
    x += Inches(3.1)

# 하단: 3.3 원천시스템 연계관리 (핵심 신규)
add_text(slide, Inches(0.8), Inches(3.35), Inches(10), Inches(0.35),
         '3.3 원천시스템 연계관리 (핵심 신규)', font_size=14, bold=True, color=RGBColor(0xC6,0x28,0x28))

conn_screens = [
    ('3.3.a CDC 연계현황 [신규]', RGBColor(0xE6,0x51,0x00), [
        'CDC서버 4대 상태 모니터링',
        'SAP/ORACLE 동기화 현황',
        '(자산/시설/수질/운영/관망통합)',
        'Lag 현황 / 지연 알림',
        '테이블별 동기화 이력',
    ]),
    ('3.3.b Kafka 스트리밍 [신규]', RGBColor(0xD3,0x2F,0x2F), [
        '카프카 클러스터(6대) 상태',
        '토픽 목록 / 파티션 현황',
        '프로듀서(4대) 상태',
        '(RWIS/HDAPS/GIOS/스마트미터링)',
        '컨슈머 그룹 Lag',
        '처리량(Throughput) 차트',
    ]),
    ('3.3.c 외부연계 [신규]', ORANGE, [
        '기상청 API 연계 상태',
        '티베로(제타7) 보정데이터',
        '호출 이력 / 성공률',
        '데이터 수신 현황',
    ]),
    ('3.3.d 단방향전송 [신규]', RGBColor(0xBF,0x36,0x0C), [
        '단방향전송장치 상태',
        'RWIS 전송 현황',
        'HDAPS 전송 현황',
        'GIOS 전송 현황',
    ]),
]

x = Inches(0.5)
for title, color, items in conn_screens:
    add_card(slide, x, Inches(3.8), Inches(3.0), Inches(2.3), title, items, title_color=color)
    x += Inches(3.1)

# 하단 추가: 3.6, 3.7
add_card(slide, Inches(0.5), Inches(6.3), Inches(3.0), Inches(0.9), '3.6 실시간 로그·알림', ['실시간 로그 스트림', '심각도 필터 (ERROR/WARN/INFO)'], title_color=RGBColor(0x78,0x90,0x9C))
add_card(slide, Inches(3.6), Inches(6.3), Inches(3.0), Inches(0.9), '3.7 정제·변환 (dbt)', ['dbt 모델 목록 / DAG 시각화', '★ ODS→분석DB 변환현황 [신규]'], title_color=RGBColor(0x78,0x90,0x9C))

add_slide_number(slide, 8, TOTAL_SLIDES)


# ======================================================================
# SLIDE 9: ④ 유통·활용 상세
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.4),
         '07  ④ 유통·활용 상세', font_size=22, bold=True, color=KWATER_DARK)
add_shape(slide, Inches(0.8), Inches(0.75), Inches(2), Inches(0.04), fill_color=ORANGE)

dist_screens = [
    ('4.1 Data Product\n구성·관리', ORANGE, [
        'Data Product 목록 (도메인별)',
        '등록 폼',
        '상세 (구성정보/SLA/이력)',
        '상태 변경 (활성/중지/제한/폐기)',
    ]),
    ('4.2 비식별화·\n가명처리', ORANGE, [
        '비식별화 규칙 관리',
        '처리 대상 데이터셋 선택',
        '처리 결과 미리보기',
        '처리 이력',
    ]),
    ('4.3 데이터\n신청·승인', ORANGE, [
        '신청 목록 / 내 신청현황',
        '신청서 작성',
        '(데이터셋/사유/기간/용도)',
        '승인 워크플로우 (다단계)',
        '이용 약관 / 동의 관리',
    ]),
    ('4.4 유통 API 관리', ORANGE, [
        'API 목록 (REST/Kafka/파일)',
        'API 문서 (Swagger/OpenAPI)',
        'API Key 발급/관리',
        '호출 통계 / Rate Limit',
        'API 테스트 콘솔',
    ]),
    ('4.5 외부시스템\n데이터 제공 [신규]', RGBColor(0xE6,0x51,0x00), [
        '★ DT 연계 API 현황',
        '  → 정수장DT / 관로DT',
        '  → 댐DT / 발전DT / 스마트DT',
        '제공 데이터셋 목록 / 포맷',
        '호출 이력 / SLA 충족률',
        '★ API서버 VM(x4) 상태',
    ]),
    ('4.6 데이터\n유통통계', ORANGE, [
        '전체 유통량 추이',
        '도메인별 / 소비자별 통계',
        'Top 데이터셋 랭킹',
        '활용 목적별 분석',
    ]),
]

x = Inches(0.5)
for title, color, items in dist_screens:
    add_card(slide, x, Inches(1.2), Inches(2.0), Inches(2.5), title, items, title_color=color)
    x += Inches(2.08)

# 4.5 강조 박스
add_shape(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(1.2), fill_color=RGBColor(0xFF,0xF3,0xE0), border_color=ORANGE)
add_text(slide, Inches(0.8), Inches(4.1), Inches(11.8), Inches(0.3),
         '4.5 외부시스템 데이터 제공 — 역할 범위 정의', font_size=13, bold=True, color=RGBColor(0xE6,0x51,0x00))
add_multiline_text(slide, Inches(0.8), Inches(4.45), Inches(11.5), Inches(0.7), [
    {'text': '우리의 롤: 디지털트윈(DT) 5종에 API 또는 데이터를 제공하는 것까지가 범위', 'size': 11, 'bold': True, 'color': KWATER_DARK},
    {'text': 'DT 시스템 자체 개발/운영은 별도 팀 담당 → 포털에서는 "어떤 데이터를 어떤 DT에 어떤 API로 제공 중인지" 모니터링 화면만 구성', 'size': 10, 'color': GRAY},
    {'text': '대상 DT: 정수장DT  |  관로DT  |  댐DT  |  발전DT  |  스마트DT  →  MCP서버 VM(x4) + API서버 VM(x4) 경유', 'size': 10, 'color': GRAY},
])

add_slide_number(slide, 9, TOTAL_SLIDES)


# ======================================================================
# SLIDE 10: ⑤ 메타데이터 관리 상세
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.4),
         '08  ⑤ 메타데이터 관리 상세', font_size=22, bold=True, color=KWATER_DARK)
add_shape(slide, Inches(0.8), Inches(0.75), Inches(2), Inches(0.04), fill_color=PURPLE)

meta_screens = [
    ('5.1 표준용어·코드관리', [
        '표준용어 사전 (CRUD)',
        '도메인별 필터',
        '승인 워크플로우',
        '(등록→검토→승인)',
        '코드 값 관리',
        '(코드그룹/코드항목)',
    ]),
    ('5.2 분류체계·태그관리', [
        '분류체계 트리 뷰',
        '태그 관리 (생성/병합/삭제)',
        '태그-데이터셋 매핑 현황',
    ]),
    ('5.3 데이터모델관리', [
        '논리 모델 목록',
        '물리 모델 목록',
        '★ ERD 시각화 [개선]',
        '  → 원천시스템별 ERD (10개 DB)',
        '모델 상세 (테이블/컬럼/관계)',
        '변경이력 관리',
    ]),
    ('5.4 DQ 검증·프로파일링', [
        '품질 규칙 관리',
        '프로파일링 실행/결과',
        '★ 원천시스템별 DQ 현황 [신규]',
        '  → 10개 DB별 품질 대시보드',
        '이상 탐지 알림',
    ]),
    ('5.5 온톨로지 관리 [신규]', [
        '★ 그래프DB 스키마 관리',
        '클래스/속성 탐색 (트리 뷰)',
        '★ 마인즈DB 메타데이터 현황',
        '온톨로지 매핑 규칙',
        '',
        '* 기존 카탈로그 ②에서 분리',
        '  → 관리 기능은 ⑤에 배치',
        '  → 탐색 기능은 ②에 유지',
    ]),
]

x = Inches(0.5)
for title, items in meta_screens:
    add_card(slide, x, Inches(1.2), Inches(2.4), Inches(3.2), title, items, title_color=PURPLE)
    x += Inches(2.5)

add_slide_number(slide, 10, TOTAL_SLIDES)


# ======================================================================
# SLIDE 11: ⑥ 시스템관리 상세
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.4),
         '09  ⑥ 시스템관리 상세', font_size=22, bold=True, color=KWATER_DARK)
add_shape(slide, Inches(0.8), Inches(0.75), Inches(2), Inches(0.04), fill_color=RGBColor(0x45,0x5A,0x64))

sys_color = RGBColor(0x45,0x5A,0x64)

sys_top = [
    ('6.1 조직 및 사용자관리', [
        '사용자 목록 (CRUD)',
        '부서/조직 트리',
        '신규 가입 승인',
        '외부사용자 관리',
    ]),
    ('6.2 권한 및 역할관리', [
        '역할 정의 매트릭스',
        '역할별 메뉴 접근 설정',
        '데이터 접근 등급 설정',
        '역할 할당 이력',
    ]),
    ('6.3 데이터등급·보안정책', [
        '4단계 등급 정의',
        '(비밀/대외비/업무/공개)',
        '등급별 정책 (암호화/접근/로깅)',
        '데이터셋-등급 매핑',
        '보안 감사 현황',
    ]),
]

x = Inches(0.5)
for title, items in sys_top:
    add_card(slide, x, Inches(1.1), Inches(3.8), Inches(2.0), title, items, title_color=sys_color)
    x += Inches(4.0)

# 6.4 인프라 현황 (핵심 신규)
add_text(slide, Inches(0.8), Inches(3.35), Inches(10), Inches(0.3),
         '6.4 인프라 현황 [완전 신규] — 아키텍처 기반 서버 모니터링', font_size=13, bold=True, color=sys_color)

infra_sections = [
    ('수집계 인프라', RGBColor(0xC6,0x28,0x28), [
        '수집DB PM(x3)',
        '  → CPU/메모리/디스크/커넥션',
        '카프카서버 VM(x6)',
        '  → 브로커 상태',
        'CDC서버 VM(x4)',
        '수집Agent VM(x5)',
    ]),
    ('분석계 인프라', RGBColor(0x00,0x96,0x88), [
        'ODSDB PM(x3)',
        '  → 상태/용량',
        '분석DB PM(x4)',
        '  → 상태/용량',
    ]),
    ('서비스계 인프라', KWATER_BLUE, [
        'WEB서버 VM(x4)',
        'WAS서버 VM(x4)',
        '모니터링서버 VM(x4)',
    ]),
]

x = Inches(0.5)
for title, color, items in infra_sections:
    add_card(slide, x, Inches(3.8), Inches(3.8), Inches(2.3), title, items, title_color=color)
    x += Inches(4.0)

# 하단: 6.5 ~ 6.7
add_card(slide, Inches(0.5), Inches(6.3), Inches(3.8), Inches(0.9), '6.5 LLM·AI 엔진관리 [신규]',
         ['LLM GM(x1) GPU B300 사용률  |  모델 버전  |  질의 이력  |  그래프DB PM(x1)'], title_color=PURPLE)
add_card(slide, Inches(4.5), Inches(6.3), Inches(3.8), Inches(0.9), '6.6 연계인터페이스 모니터링',
         ['★ MCP서버 VM(x4) [신규]  |  API서버 VM(x4)  |  DT 연계 5종  |  외부연계'], title_color=sys_color)
add_card(slide, Inches(8.5), Inches(6.3), Inches(4.3), Inches(0.9), '6.7 접속통계·감사로그',
         ['일별/월별 접속통계  |  역할별 접속  |  데이터접근 감사로그  |  이상접근 탐지'], title_color=sys_color)

add_slide_number(slide, 11, TOTAL_SLIDES)


# ======================================================================
# SLIDE 12: RBAC 매트릭스
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.4),
         '10  역할별 접근권한 매트릭스 (RBAC)', font_size=22, bold=True, color=KWATER_DARK)
add_shape(slide, Inches(0.8), Inches(0.75), Inches(2), Inches(0.04), fill_color=KWATER_BLUE)

# 범례
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.25),
         '●: 전체권한    △: 제한적(읽기)    ○: 접근불가    ◈: 본인건만', font_size=9, color=GRAY)

# 테이블 데이터
rbac_headers = ['메뉴', '수공직원\n(일반)', '부서\n관리자', '데이터\n스튜어드', '데이터\n엔지니어', '협력\n직원', '외부\n사용자', '시스템\n관리자']
rbac_data = [
    ['① 통합대시보드', '', '', '', '', '', '', ''],
    ['  1.1~1.2 홈/프로세스', '●', '●', '●', '●', '△', '△', '●'],
    ['  1.3 품질현황', '△', '●', '●', '●', '○', '○', '●'],
    ['  1.4 원천시스템 모니터링', '△', '●', '●', '●', '○', '○', '●'],
    ['  1.5 AI 검색', '●', '●', '●', '●', '△', '△', '●'],
    ['② 카탈로그·탐색', '', '', '', '', '', '', ''],
    ['  2.1~2.2 검색/상세', '●', '●', '●', '●', '△', '△', '●'],
    ['  2.3~2.4 그래프/리니지', '△', '●', '●', '●', '○', '○', '●'],
    ['③ 수집·정제', '', '', '', '', '', '', ''],
    ['  3.1~3.7 전체', '○', '△', '●', '●', '○', '○', '●'],
    ['④ 유통·활용', '', '', '', '', '', '', ''],
    ['  4.1~4.2 상품/비식별', '△', '●', '●', '△', '△', '○', '●'],
    ['  4.3 신청·승인', '●', '●', '●', '△', '●', '●', '●'],
    ['  4.5 외부시스템 제공', '○', '△', '●', '●', '○', '○', '●'],
    ['⑤ 메타데이터', '', '', '', '', '', '', ''],
    ['  5.1~5.5 전체', '△', '△', '●', '△', '○', '○', '●'],
    ['⑥ 시스템관리', '', '', '', '', '', '', ''],
    ['  6.1~6.7 전체', '○', '○', '○', '△', '○', '○', '●'],
]

# 테이블 그리기
col_widths = [Inches(2.5), Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9)]
start_x = Inches(0.8)
start_y = Inches(1.2)
row_h = Inches(0.28)

# 헤더
x = start_x
for i, (header, w) in enumerate(zip(rbac_headers, col_widths)):
    s = add_shape(slide, x, start_y, w, Inches(0.45), fill_color=KWATER_DARK)
    s.text_frame.paragraphs[0].text = header
    s.text_frame.paragraphs[0].font.size = Pt(7.5)
    s.text_frame.paragraphs[0].font.bold = True
    s.text_frame.paragraphs[0].font.color.rgb = WHITE
    s.text_frame.paragraphs[0].font.name = '맑은 고딕'
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    s.text_frame.word_wrap = True
    x += w

# 데이터 행
y = start_y + Inches(0.45)
for row in rbac_data:
    x = start_x
    is_section = row[1] == ''
    bg = KWATER_LIGHT if is_section else WHITE
    for i, (cell, w) in enumerate(zip(row, col_widths)):
        s = add_shape(slide, x, y, w, row_h, fill_color=bg, border_color=RGBColor(0xE0,0xE0,0xE0), border_width=Pt(0.5))
        s.text_frame.paragraphs[0].text = cell
        s.text_frame.paragraphs[0].font.size = Pt(8) if i == 0 else Pt(9)
        s.text_frame.paragraphs[0].font.bold = is_section and i == 0
        s.text_frame.paragraphs[0].font.color.rgb = KWATER_DARK if i == 0 else (GREEN if cell == '●' else (ORANGE if cell == '△' else (RED if cell == '○' else KWATER_BLUE)))
        s.text_frame.paragraphs[0].font.name = '맑은 고딕'
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
        x += w
    y += row_h

add_slide_number(slide, 12, TOTAL_SLIDES)


# ======================================================================
# SLIDE 13: 화면 플로우
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.4),
         '11  화면 플로우 / 네비게이션 맵', font_size=22, bold=True, color=KWATER_DARK)
add_shape(slide, Inches(0.8), Inches(0.75), Inches(2), Inches(0.04), fill_color=KWATER_BLUE)

# 로그인 → 홈
s = add_rounded_shape(slide, Inches(5.5), Inches(1.1), Inches(2.3), Inches(0.5), fill_color=RGBColor(0x26,0x32,0x38))
s.text_frame.paragraphs[0].text = '로그인 (역할 선택)'
s.text_frame.paragraphs[0].font.size = Pt(10)
s.text_frame.paragraphs[0].font.bold = True
s.text_frame.paragraphs[0].font.color.rgb = WHITE
s.text_frame.paragraphs[0].font.name = '맑은 고딕'
s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 화살표 대체 텍스트
add_text(slide, Inches(6.3), Inches(1.6), Inches(0.7), Inches(0.3), '▼', font_size=16, color=KWATER_BLUE, alignment=PP_ALIGN.CENTER)

# 홈 (중앙)
s = add_rounded_shape(slide, Inches(4.5), Inches(1.9), Inches(4.3), Inches(0.9), fill_color=KWATER_BLUE)
tf = s.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '1.1 업무중심홈 (랜딩)'
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = '맑은 고딕'
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = 'AI퀵서치 | KPI | 업무현황 | 품질요약 | DB상태'
p2.font.size = Pt(8)
p2.font.color.rgb = RGBColor(0xCC,0xDD,0xFF)
p2.font.name = '맑은 고딕'
p2.alignment = PP_ALIGN.CENTER

# 6개 대메뉴 연결
nav_flow = [
    ('① 통합대시보드', KWATER_BLUE, 0.3),
    ('② 카탈로그·탐색', RGBColor(0x00,0x96,0x88), 2.35),
    ('③ 수집·정제', RGBColor(0xC6,0x28,0x28), 4.4),
    ('④ 유통·활용', ORANGE, 6.45),
    ('⑤ 메타데이터', PURPLE, 8.5),
    ('⑥ 시스템관리', RGBColor(0x45,0x5A,0x64), 10.55),
]

for title, color, xpos in nav_flow:
    s = add_rounded_shape(slide, Inches(xpos), Inches(3.2), Inches(1.9), Inches(0.45), fill_color=color)
    s.text_frame.paragraphs[0].text = title
    s.text_frame.paragraphs[0].font.size = Pt(9)
    s.text_frame.paragraphs[0].font.bold = True
    s.text_frame.paragraphs[0].font.color.rgb = WHITE
    s.text_frame.paragraphs[0].font.name = '맑은 고딕'
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 크로스 네비게이션
add_text(slide, Inches(0.8), Inches(4.0), Inches(10), Inches(0.35),
         '주요 크로스 네비게이션 (Cross Navigation)', font_size=13, bold=True, color=KWATER_DARK)

cross_navs = [
    ('카탈로그 검색 → 상세정보', '→ 리니지 보기 / 데이터 신청(④) / 품질 보기(⑤) / 표준용어(⑤)', KWATER_BLUE),
    ('원천시스템 필터 클릭', '→ 계측DB: RWIS|HDAPS|GIOS|스마트미터링 → 1.4.1 대시보드로 이동', RGBColor(0xC6,0x28,0x28)),
    ('원천시스템 필터 클릭', '→ 자산DB: 자산|시설|수질|운영|관망통합|기상청 → 1.4.2 대시보드로 이동', KWATER_BLUE),
    ('리니지 노드 클릭', '→ 파이프라인 상세(③) / dbt 모델(③) / API 연계(④)', RGBColor(0x00,0x96,0x88)),
    ('파이프라인 관리', '→ CDC 현황(③3.3a) / Kafka 현황(③3.3b) / 외부연계(③3.3c)', RGBColor(0xC6,0x28,0x28)),
    ('Data Product 상세', '→ API 문서(④4.4) / DT 연계 현황(④4.5) / 유통통계(④4.6)', ORANGE),
    ('AI 검색 결과', '→ 데이터셋 상세(②) / 품질 대시보드(①) / 리니지(②)', PURPLE),
]

y = Inches(4.45)
for source, target, color in cross_navs:
    s = add_rounded_shape(slide, Inches(0.8), y, Inches(2.8), Inches(0.3), fill_color=color)
    s.text_frame.paragraphs[0].text = source
    s.text_frame.paragraphs[0].font.size = Pt(8)
    s.text_frame.paragraphs[0].font.bold = True
    s.text_frame.paragraphs[0].font.color.rgb = WHITE
    s.text_frame.paragraphs[0].font.name = '맑은 고딕'
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(slide, Inches(3.7), y + Inches(0.04), Inches(9), Inches(0.25), target, font_size=8.5, color=GRAY)
    y += Inches(0.38)

add_slide_number(slide, 13, TOTAL_SLIDES)


# ======================================================================
# SLIDE 14: 변경점 요약 + Next Steps
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.4),
         '12  현재 대비 변경점 요약 및 Next Steps', font_size=22, bold=True, color=KWATER_DARK)
add_shape(slide, Inches(0.8), Inches(0.75), Inches(2), Inches(0.04), fill_color=KWATER_BLUE)

# 화면 수 비교
add_text(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.3),
         '화면 수 비교', font_size=14, bold=True, color=KWATER_DARK)

compare_data = [
    ['구분', '현재', '재설계', '변경'],
    ['① 통합대시보드', '7개', '7+10', '계측4+자산6 추가'],
    ['② 카탈로그·탐색', '7개', '6개', '온톨로지→⑤이동'],
    ['③ 수집·정제', '8개', '7+4', 'CDC/Kafka/외부연계 추가'],
    ['④ 유통·활용', '5개', '6개', '외부시스템제공 추가'],
    ['⑤ 메타데이터', '4개', '5개', '온톨로지관리 추가'],
    ['⑥ 시스템관리', '6개', '7+3', '인프라/LLM 추가'],
    ['합계', '37개', '38+17', ''],
]

cy = Inches(1.4)
for i, row in enumerate(compare_data):
    cx = Inches(0.8)
    cws = [Inches(2.2), Inches(1.0), Inches(1.0), Inches(2.5)]
    for j, (cell, cw) in enumerate(zip(row, cws)):
        bg = KWATER_DARK if i == 0 else (LIGHT_GRAY if i % 2 == 0 else WHITE)
        fc = WHITE if i == 0 else KWATER_DARK
        s = add_shape(slide, cx, cy, cw, Inches(0.28), fill_color=bg, border_color=RGBColor(0xE0,0xE0,0xE0), border_width=Pt(0.5))
        s.text_frame.paragraphs[0].text = cell
        s.text_frame.paragraphs[0].font.size = Pt(9)
        s.text_frame.paragraphs[0].font.bold = i == 0
        s.text_frame.paragraphs[0].font.color.rgb = fc
        s.text_frame.paragraphs[0].font.name = '맑은 고딕'
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cx += cw
    cy += Inches(0.28)

# 주요 변경 내역
add_text(slide, Inches(0.8), Inches(3.9), Inches(6), Inches(0.3),
         '주요 변경 14건', font_size=14, bold=True, color=KWATER_DARK)

changes = [
    ('[신규] 1.4.2 자산DB 대시보드 — 6개 DB 전용 화면',
    '[개편] 1.4.1 계측DB — GIOS 추가, 명칭 통일',
    '[신규] 3.3.a CDC 연계현황 (서버 4대)',
    '[신규] 3.3.b Kafka 스트리밍 모니터링 (서버 6대)',
    '[신규] 3.3.c 외부연계 (기상청 API + 티베로)',
    '[신규] 3.3.d 단방향전송 모니터링',
    '[신규] 4.5 외부시스템 데이터 제공 (DT 5종)',
    '[분리] 5.5 온톨로지 관리 (②→⑤ 이동)',
    '[신규] 6.4 인프라 현황 (수집/분석/서비스)',
    '[신규] 6.5 LLM·AI 엔진관리 (GPU B300)',
    '[강화] 6.6 연계 모니터링 — MCP서버 4대',
    '[개선] 전체 검색 필터에 10개 원천시스템 명시',
    '[삭제] 온톨로지 클래스탐색 → 5.5 통합',
    '[삭제] dbt 모델 상세 → 3.7 통합')
]

y = Inches(4.2)
for i, change in enumerate(changes):
    color = GREEN if '[신규]' in change else (KWATER_BLUE if '[개편]' in change or '[개선]' in change or '[강화]' in change else (RED if '[삭제]' in change else PURPLE))
    add_text(slide, Inches(0.8), y, Inches(6.5), Inches(0.2), f'{i+1}. {change}', font_size=8, color=color)
    y += Inches(0.2)

# Next Steps (우측)
add_text(slide, Inches(7.8), Inches(1.0), Inches(5), Inches(0.3),
         'Next Steps', font_size=14, bold=True, color=KWATER_DARK)

steps = [
    {'text': 'Phase 1: 사이드바 메뉴 구조 변경', 'size': 11, 'bold': True, 'color': KWATER_DARK},
    {'text': '  → IA 기준 메뉴 재구성 (코드 반영)', 'size': 9, 'color': GRAY},
    {'text': '', 'size': 6},
    {'text': 'Phase 2: 신규 화면 프로토타입', 'size': 11, 'bold': True, 'color': KWATER_DARK},
    {'text': '  → 1.4.2 자산DB 대시보드', 'size': 9, 'color': GRAY},
    {'text': '  → 3.3 원천시스템 연계관리 (CDC/Kafka)', 'size': 9, 'color': GRAY},
    {'text': '  → 4.5 외부시스템 데이터 제공', 'size': 9, 'color': GRAY},
    {'text': '  → 6.4 인프라 현황', 'size': 9, 'color': GRAY},
    {'text': '  → 6.5 LLM·AI 엔진관리', 'size': 9, 'color': GRAY},
    {'text': '', 'size': 6},
    {'text': 'Phase 3: 기존 화면 개선', 'size': 11, 'bold': True, 'color': KWATER_DARK},
    {'text': '  → 검색 필터 10개 DB 반영', 'size': 9, 'color': GRAY},
    {'text': '  → 계측DB 명칭 통일 (GIOS 추가)', 'size': 9, 'color': GRAY},
    {'text': '  → 다크모드 인라인 색상 수정', 'size': 9, 'color': GRAY},
    {'text': '  → 역할별 RBAC 메뉴 필터링', 'size': 9, 'color': GRAY},
    {'text': '', 'size': 6},
    {'text': 'Phase 4: UX 고도화', 'size': 11, 'bold': True, 'color': KWATER_DARK},
    {'text': '  → 알림 시스템, 페이지네이션', 'size': 9, 'color': GRAY},
    {'text': '  → 반응형 레이아웃', 'size': 9, 'color': GRAY},
    {'text': '  → 접근성 개선', 'size': 9, 'color': GRAY},
]

add_multiline_text(slide, Inches(7.8), Inches(1.4), Inches(5), Inches(5.5), steps)

add_slide_number(slide, 14, TOTAL_SLIDES)


# ===== SAVE =====
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'DataHub_Portal_IA_재설계_v3.0.pptx')
output_path = os.path.normpath(output_path)
prs.save(output_path)
print(f'PPT 생성 완료: {output_path}')
